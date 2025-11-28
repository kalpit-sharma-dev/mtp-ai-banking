"""
Script to generate PowerPoint presentation for AI Banking Project
Requires: python-pptx library
Install: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(0, 51, 102)  # Dark blue
    subtitle_color = RGBColor(51, 51, 51)  # Dark gray
    accent_color = RGBColor(0, 102, 204)  # Blue
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Banking Using Agentic AI and ML Models"
    subtitle.text = "Involving MCP Servers\nProject Presentation"
    
    # Slide 2: Problem Definition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Problem Definition"
    
    text_frame = content.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Current Banking Systems Face:"
    p.font.bold = True
    p.font.size = Pt(14)
    
    problems = [
        "• Lack of intelligent, context-aware processing",
        "• Difficulty integrating multiple AI/ML models",
        "• Inability to scale securely across banking services",
        "• No unified orchestration layer for AI agents",
        "",
        "Specific Problems:",
        "• Traditional rule-based systems cannot adapt",
        "• ML models exist in silos without coordination",
        "• No context-aware routing between channels",
        "• Manual intervention required for decisions",
        "• Inconsistent logic across channels"
    ]
    
    for problem in problems:
        p = text_frame.add_paragraph()
        p.text = problem
        p.font.size = Pt(11)
        p.level = 0 if problem.startswith("•") else 0
    
    # Slide 3: Why This Problem Is Important
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Why This Problem Is Important"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    reasons = [
        ("1. Financial Security & Fraud Prevention", [
            "• Financial frauds increasing exponentially",
            "• Real-time intelligent defenses essential",
            "• Static rule engines cannot adapt",
            "• Billions lost annually to fraud"
        ]),
        ("2. Customer Expectations", [
            "• Users demand personalized, predictive services",
            "• Modern customers expect AI-powered experiences",
            "• 24/7 intelligent assistance required"
        ]),
        ("3. Regulatory Compliance", [
            "• Explainability in AI decisions required",
            "• Audit trails for all decisions",
            "• RBI regulations demand robust risk management"
        ]),
        ("4. Competitive Advantage", [
            "• Differentiate through AI capabilities",
            "• First-mover advantage in AI-driven banking",
            "• Future-ready infrastructure"
        ])
    ]
    
    for reason_title, points in reasons:
        p = text_frame.add_paragraph()
        p.text = reason_title
        p.font.bold = True
        p.font.size = Pt(12)
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(10)
            p.level = 1
    
    # Slide 4: What the Project Delivers
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "What the Project Delivers"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    deliverables = [
        ("1. Unified AI Orchestration Layer", [
            "• MCP Server (Golang) for centralized orchestration",
            "• AI Skin Orchestrator for intent recognition",
            "• Agent Mesh with specialized AI agents"
        ]),
        ("2. Intelligent Agent System", [
            "• Banking Agent: Transactions, balance, statements",
            "• Fraud Agent: Real-time fraud detection",
            "• Guardrail Agent: RBI regulations enforcement",
            "• Clearance Agent: Loan approvals",
            "• Scoring Agent: Credit and risk scoring"
        ]),
        ("3. ML Models Integration", [
            "• Fraud Detection: XGBoost classifier",
            "• Credit Scoring: Random Forest",
            "• Risk Scoring: Ensemble model"
        ]),
        ("4. Key Features", [
            "• Real-time decisions (<200ms latency)",
            "• Context-aware processing",
            "• Explainable AI decisions",
            "• Scalable (millions of transactions/day)"
        ])
    ]
    
    for deliverable_title, points in deliverables:
        p = text_frame.add_paragraph()
        p.text = deliverable_title
        p.font.bold = True
        p.font.size = Pt(12)
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(10)
            p.level = 1
    
    # Slide 5: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Architecture Overview"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    architecture = [
        "5-Layer Architecture:",
        "",
        "Layer 0: Web UI (Port 3000)",
        "• User interface for banking operations",
        "• Natural language input support",
        "",
        "Layer 1: MCP Server (Port 8080)",
        "• Central orchestration hub",
        "• Task routing and management",
        "• Session management",
        "",
        "Layer 2: AI Skin Orchestrator (Port 8081)",
        "• Intent recognition (LLM + Rule-based)",
        "• Context enrichment",
        "• Natural language understanding",
        "",
        "Layer 3: Agent Mesh (Ports 8001-8005)",
        "• Banking, Fraud, Guardrail, Clearance, Scoring Agents",
        "",
        "Layer 4: ML Models Service (Port 9000)",
        "• Fraud Detection, Credit Scoring, Risk Scoring",
        "",
        "Layer 5: Banking Integrations (Port 7000)",
        "• MB, NB, DWH services"
    ]
    
    for line in architecture:
        p = text_frame.add_paragraph()
        p.text = line
        if line.endswith(":") or (line and not line.startswith("•") and line != ""):
            p.font.bold = True
        p.font.size = Pt(10)
    
    # Slide 6: Project Diagram Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Project Diagram Architecture"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Add text box for diagram description
    diagram_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    diagram_frame = diagram_text.text_frame
    diagram_frame.text = """Complete System Flow:

User Requests (MB, NB, Web UI)
    ↓
AI Skin Orchestrator (Port 8081)
    • Intent Recognition
    • Context Enrichment
    ↓
MCP Server (Port 8080)
    • Task Orchestration
    • Agent Routing
    ↓
Agent Mesh (Ports 8001-8005)
    • Banking Agent
    • Fraud Agent
    • Guardrail Agent
    • Clearance Agent
    • Scoring Agent
    ↓
ML Models Service (Port 9000)
    • Fraud Detection (XGBoost)
    • Credit Scoring (Random Forest)
    • Risk Scoring (Ensemble)
    ↓
Banking Integrations (Port 7000)
    • MB, NB, DWH Services"""
    
    for paragraph in diagram_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = "Courier New"
    
    # Slide 7: What I Implemented
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "What I Implemented (Technical Contribution)"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    implementations = [
        ("1. MCP Server (Golang) - Layer 1", [
            "• RESTful API with Gorilla Mux",
            "• Task submission and management",
            "• Session management",
            "• Context routing engine",
            "• Agent registry and discovery"
        ]),
        ("2. AI Skin Orchestrator - Layer 2", [
            "• Intent parsing (3 methods)",
            "• Context enrichment service",
            "• Ollama/Llama 3 integration",
            "• Response merging"
        ]),
        ("3. Agent Mesh - Layer 3", [
            "• 5 specialized agents",
            "• ML model integration",
            "• HTTP client for ML calls",
            "• Fallback mechanisms"
        ]),
        ("4. ML Models Service - Layer 4", [
            "• FastAPI-based REST service",
            "• XGBoost fraud detection",
            "• Random Forest credit scoring",
            "• Ensemble risk scoring"
        ]),
        ("5. Banking Integrations - Layer 5", [
            "• MB, NB, DWH services",
            "• Banking Gateway",
            "• Channel-specific routing"
        ])
    ]
    
    for impl_title, points in implementations:
        p = text_frame.add_paragraph()
        p.text = impl_title
        p.font.bold = True
        p.font.size = Pt(11)
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(9)
            p.level = 1
    
    # Slide 8: Gap in the Industry Addressed
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Gap in the Industry Addressed"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    gaps = [
        "Existing Solutions & Limitations:",
        "",
        "❌ Traditional Banking: Rule-based, non-adaptive",
        "❌ FinTech: No context-awareness",
        "❌ Cloud ML Platforms: No agentic AI integration",
        "❌ AI Frameworks: Not banking-optimized",
        "",
        "Our Solution Addresses:",
        "",
        "✅ Unified Agentic AI Platform",
        "   • First banking-specific MCP server",
        "   • Complete agent mesh",
        "",
        "✅ ML Model Integration",
        "   • Seamless integration with agents",
        "   • Real-time predictions",
        "",
        "✅ Banking-Specific Features",
        "   • RBI regulation compliance",
        "   • Multi-channel support",
        "",
        "✅ Production-Ready Architecture",
        "   • Scalable microservices",
        "   • Security and authentication"
    ]
    
    for line in gaps:
        p = text_frame.add_paragraph()
        p.text = line
        if line.endswith(":") or line.startswith("✅") or line.startswith("❌"):
            p.font.bold = True
        p.font.size = Pt(10)
    
    # Slide 9: Results Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Results Summary"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    results = [
        ("1. Functional Results", [
            "• Intent Recognition: 95% accuracy",
            "• ML Model Integration: Real-time predictions",
            "• Agent Performance: All agents operational"
        ]),
        ("2. Technical Metrics", [
            "• Latency: <200ms for real-time decisions",
            "• Scalability: Millions of transactions/day",
            "• Reliability: Fallback mechanisms"
        ]),
        ("3. Architecture Achievements", [
            "• 5-Layer architecture implemented",
            "• All layers communicate via REST/gRPC",
            "• ML models called by agents",
            "• Complete end-to-end flow"
        ]),
        ("4. Security & Compliance", [
            "• API key authentication",
            "• Session management",
            "• Audit trails",
            "• RBI regulation compliance"
        ])
    ]
    
    for result_title, points in results:
        p = text_frame.add_paragraph()
        p.text = result_title
        p.font.bold = True
        p.font.size = Pt(11)
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(9)
            p.level = 1
    
    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    conclusion = [
        "What We Achieved:",
        "",
        "✅ Complete AI Banking Platform",
        "   • India's first AI-Orchestrated Banking Platform",
        "   • Unified 5-layer architecture",
        "",
        "✅ Technical Innovation",
        "   • Novel MCP server architecture",
        "   • Context-aware agent orchestration",
        "",
        "✅ Production-Ready System",
        "   • Scalable microservices",
        "   • Security and authentication",
        "",
        "Key Contributions:",
        "1. First MCP servers for BFSI sector",
        "2. Unified orchestration model",
        "3. ML models with agentic AI integration",
        "",
        "Future Enhancements:",
        "• Full LLM conversational banking",
        "• Self-learning fraud patterns",
        "• Voice and AR Banking"
    ]
    
    for line in conclusion:
        p = text_frame.add_paragraph()
        p.text = line
        if line.endswith(":") or line.startswith("✅") or line.startswith("•"):
            p.font.bold = True
        p.font.size = Pt(10)
    
    # Save presentation
    output_file = "AI_Banking_Project_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library not found.")
        print("Please install it using: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentation: {e}")

